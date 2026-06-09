from flask import Flask, request, render_template, jsonify
from parser import parse_franchise
import json, os, tempfile

app = Flask(__name__)
app.config['TEMPLATES_AUTO_RELOAD'] = True
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0

FRANCHISES = ['DeltaFranchise', 'BetaFranchise', 'GammaFranchise', 'AlphaFranchise']

AVAILABLE_MONTHS = {
    'DeltaFranchise':  ['Marzo 2026', 'Abril 2026'],
    'BetaFranchise':     ['Abril 2026'],
    'GammaFranchise':       ['Marzo 2026', 'Abril 2026'],
    'AlphaFranchise': ['Marzo 2026', 'Abril 2026'],
}

@app.route('/')
def index():
    return render_template('index.html',
                           franchises=FRANCHISES,
                           available_months=AVAILABLE_MONTHS)

@app.route('/generate', methods=['POST'])
def generate():
    # Validações
    if 'excel' not in request.files:
        return 'Arquivo Excel não enviado', 400

    file      = request.files['excel']
    franchise = request.form.get('franchise', '')
    period    = request.form.get('period', '')
    period_long = request.form.get('period_long', f'01 – 30 {period}')

    if franchise not in FRANCHISES:
        return f'Franquia inválida: {franchise}', 400

    if not file.filename.endswith('.xlsx'):
        return 'Apenas arquivos .xlsx são aceitos', 400

    # Verificar disponibilidade do mês
    if period not in AVAILABLE_MONTHS.get(franchise, []):
        # Aviso mas não bloqueia — pode ser um mês novo
        pass

    # Salvar temporariamente e parsear
    tmp = tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False)
    tmp.close()
    file.save(tmp.name)

    try:
        data = parse_franchise(tmp.name, franchise, period, period_long)
    except Exception as e:
        return f'Erro ao processar o arquivo: {str(e)}', 500
    finally:
        os.unlink(tmp.name)

    return render_template(
        'dashboard.html',
        DATA_JSON=json.dumps(data, ensure_ascii=False),
        franchise=franchise,
        period=period,
        n_stores=data['n_stores'],
    )

@app.route('/export_pdf', methods=['POST'])
def export_pdf():
    franchise = request.form.get('franchise', 'Franquicia')
    period = request.form.get('period', 'Periodo')
    data_json = request.form.get('data_json', '{}')
    n_stores = request.form.get('n_stores', 0)
    theme = request.form.get('theme', 'dark')
    
    # Render the HTML template
    html_content = render_template(
        'dashboard.html',
        DATA_JSON=data_json,
        franchise=franchise,
        period=period,
        n_stores=n_stores
    )
    
    # Inject Playwright helper script to make all tabs visible and draw them
    inject_script = """
    <script>
    window.addEventListener('DOMContentLoaded', () => {
        // Force all tabs to be visible
        const tcs = document.querySelectorAll('.tc');
        tcs.forEach(tc => {
            tc.style.display = 'block';
            tc.classList.add('active');
        });
        // Draw all charts
        const tabsList = ['resumen', 'rankings', 'areas', 'velocidad', 'kbp1', 'kbp2', 'kbp3', 'delivery', 'drivethru', 'camcheck'];
        tabsList.forEach(tId => {
            drawTab(tId);
            drawn.add(tId);
        });
    });
    </script>
    """
    html_content = html_content.replace('</body>', inject_script + '</body>')
    
    # Force pdf-mode class on body
    if theme == 'light':
        html_content = html_content.replace('<body>', '<body class="pdf-mode light-theme">')
    else:
        html_content = html_content.replace('<body>', '<body class="pdf-mode">')
    
    # Save to temp file
    from playwright.sync_api import sync_playwright
    
    # Save locally to avoid permissions issues in temp folder
    temp_dir = os.path.dirname(__file__)
    with tempfile.NamedTemporaryFile(suffix='.html', dir=temp_dir, delete=False, mode='w', encoding='utf-8') as f:
        f.write(html_content)
        temp_html_path = f.name
        
    temp_pdf_path = temp_html_path.replace('.html', '.pdf')
    
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1400, "height": 900})
            page.goto("file://" + os.path.abspath(temp_html_path).replace("\\", "/"))
            page.wait_for_timeout(2000) # wait 2s for chart rendering
            page.pdf(
                path=temp_pdf_path,
                format="A4",
                landscape=True,
                print_background=True,
                margin={"top": "0.2in", "bottom": "0.2in", "left": "0.25in", "right": "0.25in"}
            )
            browser.close()
            
        with open(temp_pdf_path, 'rb') as pdf_file:
            pdf_data = pdf_file.read()
            
        # Cleanup
        os.unlink(temp_html_path)
        os.unlink(temp_pdf_path)
        
        from flask import make_response
        response = make_response(pdf_data)
        response.headers['Content-Type'] = 'application/pdf'
        response.headers['Content-Disposition'] = f'attachment; filename=IC_{franchise}_{period.replace(" ", "_")}_REPORT.pdf'
        return response
    except Exception as e:
        if os.path.exists(temp_html_path): os.unlink(temp_html_path)
        if os.path.exists(temp_pdf_path): os.unlink(temp_pdf_path)
        return f"Error generating PDF on server: {str(e)}", 500

@app.route('/export_pptx', methods=['POST'])
def export_pptx():
    franchise = request.form.get('franchise', 'Franquicia')
    period = request.form.get('period', 'Periodo')
    data_json = request.form.get('data_json', '{}')
    n_stores = request.form.get('n_stores', 0)
    theme = request.form.get('theme', 'dark')
    
    # Render the HTML template
    html_content = render_template(
        'dashboard.html',
        DATA_JSON=data_json,
        franchise=franchise,
        period=period,
        n_stores=n_stores
    )
    
    # Force pdf-mode class on body
    if theme == 'light':
        html_content = html_content.replace('<body>', '<body class="pdf-mode light-theme">')
    else:
        html_content = html_content.replace('<body>', '<body class="pdf-mode">')
    
    # Save locally to avoid permissions issues in temp folder
    temp_dir = os.path.dirname(__file__)
    with tempfile.NamedTemporaryFile(suffix='.html', dir=temp_dir, delete=False, mode='w', encoding='utf-8') as f:
        f.write(html_content)
        temp_html_path = f.name
        
    temp_pptx_path = temp_html_path.replace('.html', '.pptx')
    
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    tabs_list = ['resumen', 'rankings', 'areas', 'velocidad', 'kbp1', 'kbp2', 'kbp3', 'delivery', 'drivethru', 'camcheck']
    screenshots = []
    
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1400, "height": 900})
            page.goto("file://" + os.path.abspath(temp_html_path).replace("\\", "/"))
            
            # Wait for load and force pdf-mode (just in case)
            page.wait_for_timeout(1000)
            page.evaluate("document.body.classList.add('pdf-mode')")
            
            # Hide topbar and tabsbar completely to prevent sticky header overlays on screenshots
            page.evaluate("const tb = document.querySelector('.topbar'); if(tb) tb.style.setProperty('display', 'none', 'important')")
            page.evaluate("const tbar = document.querySelector('.tabs-bar'); if(tbar) tbar.style.setProperty('display', 'none', 'important')")
            
            for t_idx, t_id in enumerate(tabs_list):
                # Switch tab programmatically using our JS sw function
                page.evaluate(f"sw('{t_id}', document.querySelectorAll('.tab')[{t_idx}])")
                # Wait 800ms for charts to render
                page.wait_for_timeout(800)
                
                # Bounding box / screenshot of the .main element
                main_element = page.query_selector('.main')
                img_temp = tempfile.NamedTemporaryFile(suffix='.png', dir=temp_dir, delete=False)
                img_temp.close()
                
                # Take screenshot
                main_element.screenshot(path=img_temp.name)
                screenshots.append(img_temp.name)
                
                # Add slide and place screenshot
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(img_temp.name, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
                
            browser.close()
            
        prs.save(temp_pptx_path)
        
        with open(temp_pptx_path, 'rb') as pptx_file:
            pptx_data = pptx_file.read()
            
        # Cleanup
        os.unlink(temp_html_path)
        os.unlink(temp_pptx_path)
        for path in screenshots:
            if os.path.exists(path):
                os.unlink(path)
                
        from flask import make_response
        response = make_response(pptx_data)
        response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        response.headers['Content-Disposition'] = f'attachment; filename=IC_{franchise}_{period.replace(" ", "_")}_REPORT.pptx'
        return response
        
    except Exception as e:
        if os.path.exists(temp_html_path): os.unlink(temp_html_path)
        if os.path.exists(temp_pptx_path): os.unlink(temp_pptx_path)
        for path in screenshots:
            if os.path.exists(path): os.unlink(path)
        return f"Error generating PowerPoint on server: {str(e)}", 500


@app.route('/js-error', methods=['POST'])
def js_error():
    try:
        err = request.get_json()
        log_path = os.path.join(os.path.dirname(__file__), 'client_errors.log')
        with open(log_path, 'a', encoding='utf-8') as f:
            f.write(json.dumps(err, ensure_ascii=False) + '\n')
        print("LOGGER: Client JS error written to log!")
    except Exception as e:
        print("LOGGER ERROR:", e)
    return 'OK'

if __name__ == '__main__':
    app.run(debug=True, port=5000)

